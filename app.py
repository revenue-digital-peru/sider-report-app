"""
Sider Express — Report Web App v3 (Google Gemini)
===================================================
Igual que v2 pero usa Google Gemini (gratuito) en lugar de Anthropic API.
"""

import os
import io
import json
import re
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template
import google.generativeai as genai
import openpyxl

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

# ─────────────────────────────────────────────
# RUTAS
# ─────────────────────────────────────────────
@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    try:
        if "xlsx_actual" not in request.files:
            return jsonify({"error": "Falta el XLSX de la semana actual."}), 400
        if "xlsx_anterior" not in request.files:
            return jsonify({"error": "Falta el XLSX de la semana anterior."}), 400
        if "pptx" not in request.files:
            return jsonify({"error": "Falta el archivo PPTX plantilla."}), 400

        xlsx_actual   = request.files["xlsx_actual"].read()
        xlsx_anterior = request.files["xlsx_anterior"].read()
        pptx_bytes    = request.files["pptx"].read()

        config = {
            "mes":            request.form.get("mes", ""),
            "periodo_inicio": request.form.get("periodo_inicio", ""),
            "periodo_fin":    request.form.get("periodo_fin", ""),
            "semana":         request.form.get("semana", ""),
        }

        data_actual   = extract_xlsx_data(xlsx_actual)
        data_anterior = extract_xlsx_data(xlsx_anterior)
        variaciones   = calcular_variaciones(data_actual, data_anterior)
        texts         = generate_texts_with_gemini(data_actual, data_anterior, variaciones, config)
        replacements  = build_replacements(data_actual, data_anterior, variaciones, texts, config)
        output_bytes  = fill_pptx(pptx_bytes, replacements)

        mes      = config.get("mes", "reporte").upper()
        pi       = config.get("periodo_inicio", "").replace("/", "")
        pf       = config.get("periodo_fin", "").replace("/", "")
        filename = f"Sider_Reporte_{mes}_{pi}_{pf}.pptx"

        return send_file(
            io.BytesIO(output_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


# ─────────────────────────────────────────────
# EXTRACCIÓN DE DATOS DEL XLSX
# ─────────────────────────────────────────────
def extract_xlsx_data(xlsx_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
    sheets = wb.sheetnames

    def read_sheet(name, max_row=30):
        if name not in sheets:
            return []
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(min_row=1, max_row=max_row, values_only=True):
            if any(v is not None for v in row):
                rows.append(list(row))
        return rows

    def find_header_row(rows):
        for i, row in enumerate(rows):
            if row and any("plataforma" in str(v).lower() for v in row if v):
                return i
        return 0

    def rows_to_dicts(rows, header_idx=0):
        if not rows or len(rows) <= header_idx:
            return []
        headers = [str(h).strip() if h is not None else f"col{i}"
                   for i, h in enumerate(rows[header_idx])]
        result = []
        for row in rows[header_idx + 1:]:
            if any(v is not None for v in row):
                padded = list(row) + [None] * (len(headers) - len(row))
                result.append(dict(zip(headers, padded)))
        return result

    pres_rows     = read_sheet("PRESUPUESTO", 15)
    total_rows    = read_sheet("TOTAL", 30)
    reg_comp_rows = read_sheet("CAMPAÑAS REGULARES COMPARATIVA ", 20)
    reg_met_rows  = read_sheet("CAMPAÑAS REGULARES METRICAS SEC", 20)
    inc_comp_rows = read_sheet("CAMPAÑAS INCENTIVOS COMPARATIVA", 15)
    inc_met_rows  = read_sheet("CAMPAÑAS INCENTIVOS METRICAS SE", 15)
    funnel_rows   = read_sheet("FUNNEL (REPORTE MENSUAL) ", 15)

    return {
        "pres":      rows_to_dicts(pres_rows,     find_header_row(pres_rows)),
        "total_raw": total_rows,
        "reg_comp":  rows_to_dicts(reg_comp_rows, find_header_row(reg_comp_rows)),
        "reg_met":   rows_to_dicts(reg_met_rows,  find_header_row(reg_met_rows)),
        "inc_comp":  rows_to_dicts(inc_comp_rows, find_header_row(inc_comp_rows)),
        "inc_met":   rows_to_dicts(inc_met_rows,  find_header_row(inc_met_rows)),
        "funnel":    funnel_rows,
    }


# ─────────────────────────────────────────────
# CÁLCULO AUTOMÁTICO DE VARIACIONES
# ─────────────────────────────────────────────
def to_float(val):
    if val is None:
        return None
    try:
        cleaned = str(val).replace(",", "").replace("$", "").replace("%", "").replace(" ", "").strip()
        return float(cleaned)
    except (ValueError, TypeError):
        return None


def calc_var(actual, anterior):
    a = to_float(actual)
    b = to_float(anterior)
    if a is None or b is None or b == 0:
        return "-"
    var = ((a - b) / abs(b)) * 100
    sign = "+" if var >= 0 else ""
    return f"{sign}{var:.2f}%"


def fmt_currency(val):
    v = to_float(val)
    if v is None:
        return str(val) if val else "-"
    return f"${v:.2f}"


def _find_row(rows, key_col, key_val):
    for row in rows:
        cell = str(row.get(key_col, "") or "").strip().lower()
        if key_val.lower() in cell:
            return row
    return {}


def _get_raw(rows, label, col_idx):
    for row in rows:
        if not row:
            continue
        for i in [0, 1]:
            if len(row) > i and str(row[i] or "").strip().lower() == label.lower():
                return row[col_idx] if len(row) > col_idx else None
    return None


def calcular_variaciones(actual, anterior):
    var = {}

    # Meta / Facebook totales
    row_act = _find_row(actual["reg_comp"],   "Plataforma", "Total")
    row_ant = _find_row(anterior["reg_comp"], "Plataforma", "Total")
    vals_act = list(row_act.values())
    vals_ant = list(row_ant.values())

    inv_act  = vals_act[4] if len(vals_act) > 4 else None
    inv_ant  = vals_ant[4] if len(vals_ant) > 4 else None
    lead_act = vals_act[2] if len(vals_act) > 2 else None
    lead_ant = vals_ant[2] if len(vals_ant) > 2 else None
    cpl_act  = vals_act[3] if len(vals_act) > 3 else None
    cpl_ant  = vals_ant[3] if len(vals_ant) > 3 else None

    var["fb_inv_act"]   = fmt_currency(inv_act)
    var["fb_inv_ant"]   = fmt_currency(inv_ant)
    var["fb_inv_var"]   = calc_var(inv_act, inv_ant)
    var["fb_leads_act"] = str(int(to_float(lead_act))) if to_float(lead_act) else str(lead_act or "-")
    var["fb_leads_ant"] = str(int(to_float(lead_ant))) if to_float(lead_ant) else str(lead_ant or "-")
    var["fb_leads_var"] = calc_var(lead_act, lead_ant)
    var["fb_cpl_act"]   = fmt_currency(cpl_act)
    var["fb_cpl_ant"]   = fmt_currency(cpl_ant)
    var["fb_cpl_var"]   = calc_var(cpl_act, cpl_ant)

    # Google Ads
    row_g_act  = _find_row(actual["reg_comp"],   "Plataforma", "Google")
    row_g_ant  = _find_row(anterior["reg_comp"], "Plataforma", "Google")
    vals_g_act = list(row_g_act.values())
    vals_g_ant = list(row_g_ant.values())

    g_inv_act  = vals_g_act[4] if len(vals_g_act) > 4 else None
    g_inv_ant  = vals_g_ant[4] if len(vals_g_ant) > 4 else None
    g_lead_act = vals_g_act[2] if len(vals_g_act) > 2 else None
    g_lead_ant = vals_g_ant[2] if len(vals_g_ant) > 2 else None
    g_cpl_act  = vals_g_act[3] if len(vals_g_act) > 3 else None
    g_cpl_ant  = vals_g_ant[3] if len(vals_g_ant) > 3 else None

    var["g_inv_act"]   = fmt_currency(g_inv_act)
    var["g_inv_ant"]   = fmt_currency(g_inv_ant)
    var["g_inv_var"]   = calc_var(g_inv_act, g_inv_ant)
    var["g_leads_act"] = str(int(to_float(g_lead_act))) if to_float(g_lead_act) else str(g_lead_act or "-")
    var["g_leads_ant"] = str(int(to_float(g_lead_ant))) if to_float(g_lead_ant) else str(g_lead_ant or "-")
    var["g_leads_var"] = calc_var(g_lead_act, g_lead_ant)
    var["g_cpl_act"]   = fmt_currency(g_cpl_act)
    var["g_cpl_ant"]   = fmt_currency(g_cpl_ant)
    var["g_cpl_var"]   = calc_var(g_cpl_act, g_cpl_ant)

    return var


# ─────────────────────────────────────────────
# GENERACIÓN DE TEXTOS CON GOOGLE GEMINI
# ─────────────────────────────────────────────
def generate_texts_with_gemini(data_actual, data_anterior, variaciones, config):
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel("gemini-1.5-flash")

    reg_comp_clean = [r for r in data_actual["reg_comp"]
                      if r.get("Plataforma") or r.get("Campaña")]
    inc_comp_clean = [r for r in data_actual["inc_comp"]
                      if r.get("Plataforma") or r.get("Campaña")]

    prompt = f"""
Eres un Programmatic Analyst Senior. Analiza los datos de performance de Sider Express
y genera textos analíticos ejecutivos para el reporte semanal.

PERIODO ACTUAL: {config.get('periodo_inicio')} - {config.get('periodo_fin')}
MES: {config.get('mes')}

=== CAMPAÑAS REGULARES — SEMANA ACTUAL ===
{json.dumps(reg_comp_clean, ensure_ascii=False, indent=2)}

=== CAMPAÑAS REGULARES — MÉTRICAS SECUNDARIAS ===
{json.dumps(data_actual['reg_met'], ensure_ascii=False, indent=2)}

=== CAMPAÑAS INCENTIVOS — SEMANA ACTUAL ===
{json.dumps(inc_comp_clean, ensure_ascii=False, indent=2)}

=== CAMPAÑAS INCENTIVOS — MÉTRICAS SECUNDARIAS ===
{json.dumps(data_actual['inc_met'], ensure_ascii=False, indent=2)}

=== VARIACIONES VS SEMANA ANTERIOR ===
Facebook - Inversión: {variaciones['fb_inv_ant']} → {variaciones['fb_inv_act']} ({variaciones['fb_inv_var']})
Facebook - Leads: {variaciones['fb_leads_ant']} → {variaciones['fb_leads_act']} ({variaciones['fb_leads_var']})
Facebook - CPL: {variaciones['fb_cpl_ant']} → {variaciones['fb_cpl_act']} ({variaciones['fb_cpl_var']})
Google - Inversión: {variaciones['g_inv_ant']} → {variaciones['g_inv_act']} ({variaciones['g_inv_var']})
Google - Leads: {variaciones['g_leads_ant']} → {variaciones['g_leads_act']} ({variaciones['g_leads_var']})
Google - CPL: {variaciones['g_cpl_ant']} → {variaciones['g_cpl_act']} ({variaciones['g_cpl_var']})

Genera textos analíticos ejecutivos, directos, basados 100% en los datos.
Máximo 4-5 oraciones por campo. Sin viñetas.

Responde ÚNICAMENTE con un JSON válido. Sin backticks, sin texto antes ni después.

{{
  "comentario_overview": "...",
  "comentario_metricas_sec": "...",
  "comentario_incentivos": "...",
  "comentario_inc_metricas": "...",
  "fb_comentario_cpl": "CPL (Costo por lead): ...",
  "fb_comentario_leads": "Leads: ...",
  "g_comentario_cpl": "CPL (Costo por lead): ...",
  "g_comentario_leads": "Leads: ...",
  "next_steps_meta": "bullet 1\\nbullet 2\\nbullet 3\\nbullet 4",
  "next_steps_google": "bullet 1\\nbullet 2\\nbullet 3",
  "next_steps_seguimiento": "bullet 1\\nbullet 2\\nbullet 3"
}}
"""

    response = model.generate_content(prompt)
    raw = response.text.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    return json.loads(raw.strip())


# ─────────────────────────────────────────────
# CONSTRUCCIÓN DE REEMPLAZOS
# ─────────────────────────────────────────────
def _v(val):
    if val is None:
        return "-"
    if isinstance(val, float):
        return str(int(val)) if val == int(val) else f"{val:.2f}"
    return str(val)


def build_replacements(data_actual, data_anterior, variaciones, texts, config):
    rep = {}
    act = data_actual

    rep["{{MES}}"]            = config.get("mes", "")
    rep["{{PERIODO_INICIO}}"] = config.get("periodo_inicio", "")
    rep["{{PERIODO_FIN}}"]    = config.get("periodo_fin", "")
    rep["{{SEMANA}}"]         = config.get("semana", "")

    # PRESUPUESTO
    pres_items = [
        ("Meta Perfo",    "PRES_META_PERFO"),
        ("Meta Branding", "PRES_META_BRAND"),
        ("Google",        "PRES_GOOGLE"),
        ("Tik Tok",       "PRES_TIKTOK"),
        ("Total",         "PRES_TOTAL"),
    ]
    for label, prefix in pres_items:
        row = _find_row(act["pres"], "Plataforma", label)
        rep[f"{{{{{prefix}_REAL}}}}"] = _v(row.get("Real", "-"))
        rep[f"{{{{{prefix}_META}}}}"] = _v(row.get("Meta", "-"))
        rep[f"{{{{{prefix}_LOG}}}}"]  = _v(row.get("Logrado", "-"))

    # TOTAL
    total = act["total_raw"]
    rep["{{LEADS_FB_REAL}}"]      = _v(_get_raw(total, "Facebook", 2))
    rep["{{LEADS_FB_META}}"]      = _v(_get_raw(total, "Facebook", 3))
    rep["{{LEADS_FB_LOG}}"]       = _v(_get_raw(total, "Facebook", 4))
    rep["{{LEADS_GOOGLE_REAL}}"]  = _v(_get_raw(total, "Google ", 2))
    rep["{{LEADS_TOTAL_REAL}}"]   = _v(_get_raw(total, "Total", 2))
    rep["{{LEADS_CAL_PCT}}"]      = _v(_get_raw(total, "Leads calificados", 2))
    rep["{{LEADS_CAL_REAL}}"]     = _v(_get_raw(total, "Leads calificados", 3))
    rep["{{LEADS_CAL_META}}"]     = _v(_get_raw(total, "Leads calificados", 4))
    rep["{{LEADS_CAL_LOG}}"]      = _v(_get_raw(total, "Leads calificados", 5))
    rep["{{FACT_REAL}}"]          = _v(_get_raw(total, "Facturación", 8))
    rep["{{FACT_META}}"]          = _v(_get_raw(total, "Facturación", 9))
    rep["{{FACT_LOG}}"]           = _v(_get_raw(total, "Facturación", 10))
    rep["{{VENTAS_REAL}}"]        = _v(_get_raw(total, "VENTAS", 2))
    rep["{{VENTAS_META}}"]        = _v(_get_raw(total, "VENTAS", 3))
    rep["{{VENTAS_LOG}}"]         = _v(_get_raw(total, "VENTAS", 4))

    # CAMPAÑAS REGULARES
    camp_reg = [("Lima","LIMA"),("Trujillo B2C","TRUJ_B2C"),("Trujillo B2B","TRUJ_B2B")]
    for camp, prefix in camp_reg:
        row  = _find_row(act["reg_comp"], "Campaña", camp)
        vals = list(row.values())
        rep[f"{{{{{prefix}_LEADS_ACT}}}}"] = _v(vals[2]) if len(vals) > 2 else "-"
        rep[f"{{{{{prefix}_CPL_ACT}}}}"]   = _v(vals[3]) if len(vals) > 3 else "-"
        rep[f"{{{{{prefix}_INV_ACT}}}}"]   = _v(vals[4]) if len(vals) > 4 else "-"
        rep[f"{{{{{prefix}_LEADS_ANT}}}}"] = _v(vals[5]) if len(vals) > 5 else "-"
        rep[f"{{{{{prefix}_CPL_ANT}}}}"]   = _v(vals[6]) if len(vals) > 6 else "-"
        rep[f"{{{{{prefix}_INV_ANT}}}}"]   = _v(vals[7]) if len(vals) > 7 else "-"

    row_g  = _find_row(act["reg_comp"], "Plataforma", "Google")
    vals_g = list(row_g.values())
    rep["{{GOOGLE_LEADS_ACT}}"] = _v(vals_g[2]) if len(vals_g) > 2 else "-"
    rep["{{GOOGLE_CPL_ACT}}"]   = _v(vals_g[3]) if len(vals_g) > 3 else "-"
    rep["{{GOOGLE_INV_ACT}}"]   = _v(vals_g[4]) if len(vals_g) > 4 else "-"
    rep["{{GOOGLE_LEADS_ANT}}"] = _v(vals_g[5]) if len(vals_g) > 5 else "-"
    rep["{{GOOGLE_CPL_ANT}}"]   = _v(vals_g[6]) if len(vals_g) > 6 else "-"
    rep["{{GOOGLE_INV_ANT}}"]   = _v(vals_g[7]) if len(vals_g) > 7 else "-"

    row_t  = _find_row(act["reg_comp"], "Plataforma", "Total")
    vals_t = list(row_t.values())
    rep["{{TOTAL_LEADS_ACT}}"] = _v(vals_t[2]) if len(vals_t) > 2 else "-"
    rep["{{TOTAL_CPL_ACT}}"]   = _v(vals_t[3]) if len(vals_t) > 3 else "-"
    rep["{{TOTAL_INV_ACT}}"]   = _v(vals_t[4]) if len(vals_t) > 4 else "-"
    rep["{{TOTAL_LEADS_ANT}}"] = _v(vals_t[5]) if len(vals_t) > 5 else "-"
    rep["{{TOTAL_CPL_ANT}}"]   = _v(vals_t[6]) if len(vals_t) > 6 else "-"
    rep["{{TOTAL_INV_ANT}}"]   = _v(vals_t[7]) if len(vals_t) > 7 else "-"

    # MÉTRICAS SECUNDARIAS
    met_camps = [("Lima","LIMA"),("Trujillo B2C","TRUJ_B2C"),("Trujillo B2B","TRUJ_B2B")]
    for camp, prefix in met_camps:
        row = _find_row(act["reg_met"], "Campaña", camp)
        rep[f"{{{{{prefix}_CTR}}}}"]      = _v(row.get("CTR ", row.get("CTR", "-")))
        rep[f"{{{{{prefix}_CTR_UNI}}}}"]  = _v(row.get("CTR unico", "-"))
        rep[f"{{{{{prefix}_PCT_MSJS}}}}"] = _v(row.get("% Mensajes", "-"))
        rep[f"{{{{{prefix}_CPM}}}}"]      = _v(row.get("CPM", "-"))
        rep[f"{{{{{prefix}_CLICS}}}}"]    = _v(row.get("Clic enlace", "-"))
        rep[f"{{{{{prefix}_ALCANCE}}}}"]  = _v(row.get("Alcance", "-"))

    row_tt = (_find_row(act["reg_met"], "Plataforma", "Total") or
              _find_row(act["reg_met"], "Campaña", "Total"))
    rep["{{TOTAL_CTR}}"]         = _v(row_tt.get("CTR ", row_tt.get("CTR", "-")))
    rep["{{TOTAL_CTR_UNI}}"]     = _v(row_tt.get("CTR unico", "-"))
    rep["{{TOTAL_PCT_MSJS}}"]    = _v(row_tt.get("% Mensajes", "-"))
    rep["{{TOTAL_CPM}}"]         = _v(row_tt.get("CPM", "-"))
    rep["{{TOTAL_IMPRESIONES}}"] = _v(row_tt.get("Impresiones", "-"))
    rep["{{TOTAL_ALCANCE}}"]     = _v(row_tt.get("Alcance", "-"))

    # INCENTIVOS
    camp_inc = [("Incentivos trujillo","INC_TRUJ"),("Incentivos Lima","INC_LIMA")]
    for camp, prefix in camp_inc:
        row  = _find_row(act["inc_comp"], "Campaña", camp)
        vals = list(row.values())
        rep[f"{{{{{prefix}_LEADS_ACT}}}}"] = _v(vals[2]) if len(vals) > 2 else "-"
        rep[f"{{{{{prefix}_CPL_ACT}}}}"]   = _v(vals[3]) if len(vals) > 3 else "-"
        rep[f"{{{{{prefix}_INV_ACT}}}}"]   = _v(vals[4]) if len(vals) > 4 else "-"
        rep[f"{{{{{prefix}_LEADS_ANT}}}}"] = _v(vals[5]) if len(vals) > 5 else "-"
        rep[f"{{{{{prefix}_CPL_ANT}}}}"]   = _v(vals[6]) if len(vals) > 6 else "-"
        rep[f"{{{{{prefix}_INV_ANT}}}}"]   = _v(vals[7]) if len(vals) > 7 else "-"

    row_it  = _find_row(act["inc_comp"], "Plataforma", "Total")
    vals_it = list(row_it.values())
    rep["{{INC_TOTAL_LEADS_ACT}}"] = _v(vals_it[2]) if len(vals_it) > 2 else "-"
    rep["{{INC_TOTAL_CPL_ACT}}"]   = _v(vals_it[3]) if len(vals_it) > 3 else "-"
    rep["{{INC_TOTAL_INV_ACT}}"]   = _v(vals_it[4]) if len(vals_it) > 4 else "-"
    rep["{{INC_TOTAL_LEADS_ANT}}"] = _v(vals_it[5]) if len(vals_it) > 5 else "-"
    rep["{{INC_TOTAL_CPL_ANT}}"]   = _v(vals_it[6]) if len(vals_it) > 6 else "-"
    rep["{{INC_TOTAL_INV_ANT}}"]   = _v(vals_it[7]) if len(vals_it) > 7 else "-"

    for camp, prefix in [("Incentivos trujillo","INC_TRUJ"),("Incentivos Lima","INC_LIMA")]:
        row = _find_row(act["inc_met"], "Campaña", camp)
        rep[f"{{{{{prefix}_CTR}}}}"]      = _v(row.get("CTR ", row.get("CTR", "-")))
        rep[f"{{{{{prefix}_CTR_UNI}}}}"]  = _v(row.get("CTR unico", "-"))
        rep[f"{{{{{prefix}_PCT_MSJS}}}}"] = _v(row.get("% Mensajes", "-"))
        rep[f"{{{{{prefix}_CPM}}}}"]      = _v(row.get("CPM", "-"))

    row_imt = (_find_row(act["inc_met"], "Plataforma", "Total") or
               _find_row(act["inc_met"], "Campaña", "Total"))
    rep["{{INC_TOTAL_CTR}}"]      = _v(row_imt.get("CTR ", row_imt.get("CTR", "-")))
    rep["{{INC_TOTAL_CTR_UNI}}"]  = _v(row_imt.get("CTR unico", "-"))
    rep["{{INC_TOTAL_PCT_MSJS}}"] = _v(row_imt.get("% Mensajes", "-"))
    rep["{{INC_TOTAL_CPM}}"]      = _v(row_imt.get("CPM", "-"))

    # VARIACIONES (calculadas automáticamente)
    rep["{{FB_INV_ACT}}"]   = variaciones["fb_inv_act"]
    rep["{{FB_INV_ANT}}"]   = variaciones["fb_inv_ant"]
    rep["{{FB_INV_VAR}}"]   = variaciones["fb_inv_var"]
    rep["{{FB_LEADS_ACT}}"] = variaciones["fb_leads_act"]
    rep["{{FB_LEADS_ANT}}"] = variaciones["fb_leads_ant"]
    rep["{{FB_LEADS_VAR}}"] = variaciones["fb_leads_var"]
    rep["{{FB_CPL_ACT}}"]   = variaciones["fb_cpl_act"]
    rep["{{FB_CPL_ANT}}"]   = variaciones["fb_cpl_ant"]
    rep["{{FB_CPL_VAR}}"]   = variaciones["fb_cpl_var"]
    rep["{{G_INV_ACT}}"]    = variaciones["g_inv_act"]
    rep["{{G_INV_ANT}}"]    = variaciones["g_inv_ant"]
    rep["{{G_INV_VAR}}"]    = variaciones["g_inv_var"]
    rep["{{G_LEADS_ACT}}"]  = variaciones["g_leads_act"]
    rep["{{G_LEADS_ANT}}"]  = variaciones["g_leads_ant"]
    rep["{{G_LEADS_VAR}}"]  = variaciones["g_leads_var"]
    rep["{{G_CPL_ACT}}"]    = variaciones["g_cpl_act"]
    rep["{{G_CPL_ANT}}"]    = variaciones["g_cpl_ant"]
    rep["{{G_CPL_VAR}}"]    = variaciones["g_cpl_var"]

    # TEXTOS DE GEMINI
    rep["{{COMENTARIO_OVERVIEW}}"]     = texts.get("comentario_overview", "")
    rep["{{COMENTARIO_METRICAS_SEC}}"] = texts.get("comentario_metricas_sec", "")
    rep["{{COMENTARIO_INCENTIVOS}}"]   = texts.get("comentario_incentivos", "")
    rep["{{COMENTARIO_INC_METRICAS}}"] = texts.get("comentario_inc_metricas", "")
    rep["{{FB_COMENTARIO_CPL}}"]       = texts.get("fb_comentario_cpl", "")
    rep["{{FB_COMENTARIO_LEADS}}"]     = texts.get("fb_comentario_leads", "")
    rep["{{G_COMENTARIO_CPL}}"]        = texts.get("g_comentario_cpl", "")
    rep["{{G_COMENTARIO_LEADS}}"]      = texts.get("g_comentario_leads", "")
    rep["{{NEXT_STEPS_META}}"]         = texts.get("next_steps_meta", "")
    rep["{{NEXT_STEPS_GOOGLE}}"]       = texts.get("next_steps_google", "")
    rep["{{NEXT_STEPS_SEGUIMIENTO}}"]  = texts.get("next_steps_seguimiento", "")

    return rep


# ─────────────────────────────────────────────
# RELLENO DEL PPTX
# ─────────────────────────────────────────────
def fill_pptx(pptx_bytes, replacements):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if not any(ph in full_text for ph in replacements):
                    continue
                new_text = full_text
                for placeholder, value in replacements.items():
                    new_text = new_text.replace(placeholder, str(value))
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# ─────────────────────────────────────────────
# HEALTH CHECK
# ─────────────────────────────────────────────
@app.route("/health")
def health():
    return jsonify({"status": "ok", "time": datetime.utcnow().isoformat()})


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
